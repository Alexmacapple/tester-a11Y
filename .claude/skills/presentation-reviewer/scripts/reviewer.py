#!/usr/bin/env python3
"""
Ajoute des commentaires de révision directement dans un fichier PowerPoint.

Usage:
    python reviewer.py <fichier.pptx> <rapport.json> --output <fichier_revu.pptx>
"""

import sys
import json
import argparse
from pathlib import Path
from typing import Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Erreur: Le module python-pptx n'est pas installé.", file=sys.stderr)
    print("Installez-le avec: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


class PresentationReviewer:
    """Ajoute des commentaires de révision dans une présentation PowerPoint."""

    SEVERITY_COLORS = {
        "high": RGBColor(220, 53, 69),      # Rouge
        "medium": RGBColor(255, 193, 7),    # Jaune
        "low": RGBColor(40, 167, 69)        # Vert
    }

    SEVERITY_LABELS = {
        "high": "CRITIQUE",
        "medium": "ATTENTION",
        "low": "SUGGESTION"
    }

    def __init__(self, pptx_path: str, analysis_path: str):
        """Initialise le reviewer avec le fichier PowerPoint et l'analyse JSON."""
        self.pptx_path = Path(pptx_path)
        self.analysis_path = Path(analysis_path)

        if not self.pptx_path.exists():
            raise FileNotFoundError(f"Fichier introuvable: {pptx_path}")
        if not self.analysis_path.exists():
            raise FileNotFoundError(f"Fichier d'analyse introuvable: {analysis_path}")

        self.presentation = Presentation(str(self.pptx_path))

        with open(self.analysis_path, 'r', encoding='utf-8') as f:
            self.analysis = json.load(f)

    def add_comments(self) -> None:
        """Ajoute des commentaires de révision dans la présentation."""
        print(f"📝 Ajout des commentaires dans: {self.pptx_path.name}")

        # Ajouter une slide de synthèse en début de présentation
        self._add_summary_slide()

        # Ajouter des commentaires sur chaque slide concernée
        for slide_analysis in self.analysis["slides"]:
            if slide_analysis["issues"]:
                self._add_slide_comments(slide_analysis)

        print(f"✅ {len([s for s in self.analysis['slides'] if s['issues']])} slides commentées")

    def _add_summary_slide(self) -> None:
        """Ajoute une slide de synthèse des problèmes détectés."""
        # Ajouter une nouvelle slide au début (après la page de titre)
        slide_layout = self.presentation.slide_layouts[1]  # Layout "Title and Content"
        slide = self.presentation.slides.add_slide(slide_layout)

        # Réordonner pour mettre après la première slide
        xml_slides = self.presentation.slides._sldIdLst
        slides_list = list(xml_slides)
        xml_slides.remove(slides_list[-1])
        xml_slides.insert(1, slides_list[-1])

        # Titre
        title = slide.shapes.title
        title.text = "📋 Rapport de Révision - Synthèse"

        # Contenu
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        summary = self.analysis["summary"]

        # Statistiques
        p = tf.paragraphs[0]
        p.text = f"Statistiques globales"
        p.font.bold = True
        p.font.size = Pt(18)

        stats = [
            f"Total de slides: {self.analysis['total_slides']}",
            f"Mots par slide (moyenne): {summary['avg_words_per_slide']}",
            f"Problèmes détectés: {summary['total_issues']} dont {summary['high_severity_issues']} critiques"
        ]

        for stat in stats:
            p = tf.add_paragraph()
            p.text = stat
            p.level = 1
            p.font.size = Pt(14)

        # Problèmes globaux
        if self.analysis["global_issues"]:
            p = tf.add_paragraph()
            p.text = "\nProblèmes globaux identifiés"
            p.font.bold = True
            p.font.size = Pt(18)
            p.level = 0

            for issue in self.analysis["global_issues"][:3]:  # Top 3
                p = tf.add_paragraph()
                emoji = "🔴" if issue["severity"] == "high" else "🟡" if issue["severity"] == "medium" else "🟢"
                p.text = f"{emoji} {issue['message']}"
                p.level = 1
                p.font.size = Pt(14)

    def _add_slide_comments(self, slide_analysis: Dict[str, Any]) -> None:
        """Ajoute des commentaires sur une slide spécifique."""
        slide_idx = slide_analysis["slide_number"]

        # +1 car on a ajouté la slide de synthèse
        slide = self.presentation.slides[slide_idx]

        # Calculer la position pour les commentaires (coin supérieur droit)
        left = Inches(8.5)
        top = Inches(0.3)
        width = Inches(3.0)

        for i, issue in enumerate(slide_analysis["issues"]):
            height = Inches(0.6)

            # Créer une forme pour le commentaire
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                top + (i * Inches(0.7)),
                width,
                height
            )

            # Couleur selon la sévérité
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.SEVERITY_COLORS[issue["severity"]]

            # Texte du commentaire
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)

            p = text_frame.paragraphs[0]
            p.text = f"{self.SEVERITY_LABELS[issue['severity']]}"
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)

            p = text_frame.add_paragraph()
            p.text = issue['message']
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(255, 255, 255)

    def save(self, output_path: str) -> None:
        """Sauvegarde la présentation avec les commentaires."""
        self.presentation.save(output_path)
        print(f"💾 Présentation sauvegardée: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Ajoute des commentaires de révision dans une présentation PowerPoint"
    )
    parser.add_argument("pptx_file", help="Fichier PowerPoint original")
    parser.add_argument("analysis_file", help="Fichier JSON d'analyse")
    parser.add_argument("--output", "-o", required=True, help="Fichier PowerPoint de sortie")

    args = parser.parse_args()

    try:
        reviewer = PresentationReviewer(args.pptx_file, args.analysis_file)
        reviewer.add_comments()
        reviewer.save(args.output)

        print("\n✅ Révision terminée avec succès!")

    except FileNotFoundError as e:
        print(f"❌ Erreur: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"❌ Erreur inattendue: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
