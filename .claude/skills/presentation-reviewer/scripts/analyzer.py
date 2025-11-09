#!/usr/bin/env python3
"""
Analyseur de présentations PowerPoint pour détecter les problèmes de structure,
clarté et conformité aux standards de conseil.

Usage:
    python analyzer.py <fichier.pptx> [--output rapport.json]
"""

import sys
import json
import argparse
from pathlib import Path
from typing import Dict, List, Any
from collections import Counter

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("Erreur: Le module python-pptx n'est pas installé.", file=sys.stderr)
    print("Installez-le avec: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# Import de la configuration
try:
    from config import load_config
except ImportError:
    print("⚠️  Module config.py non trouvé, utilisation de valeurs par défaut", file=sys.stderr)
    def load_config(config_file=None, preset=None):
        return {
            'max_words_per_slide': 50,
            'max_bullets_per_slide': 7,
            'min_words_per_slide': 5,
            'min_font_size': 14,
            'max_font_variations': 4,
            'dense_slide_threshold': 40,
            'max_dense_slides_consecutive': 3
        }


class PresentationAnalyzer:
    """Analyse une présentation PowerPoint selon les critères de conseil."""

    def __init__(self, pptx_path: str, config_file: str = None, preset: str = None):
        """
        Initialise l'analyseur avec le chemin du fichier PowerPoint.

        Args:
            pptx_path: Chemin vers le fichier PowerPoint
            config_file: Chemin vers un fichier de configuration JSON (optionnel)
            preset: Nom d'un preset de configuration (optionnel)
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"Fichier introuvable: {pptx_path}")

        # Charger la configuration
        self.config = load_config(config_file=config_file, preset=preset)

        self.presentation = Presentation(str(self.pptx_path))
        self.analysis = {
            "filename": self.pptx_path.name,
            "total_slides": len(self.presentation.slides),
            "config_used": {
                "file": config_file if config_file else None,
                "preset": preset if preset else "default",
                "max_words": self.config.get('max_words_per_slide', 50),
                "max_bullets": self.config.get('max_bullets_per_slide', 7)
            },
            "slides": [],
            "global_issues": [],
            "summary": {}
        }

    def analyze(self) -> Dict[str, Any]:
        """Effectue l'analyse complète de la présentation."""
        print(f"🔍 Analyse de: {self.pptx_path.name}")
        print(f"   {self.analysis['total_slides']} slides détectées\n")

        # Analyse slide par slide
        for idx, slide in enumerate(self.presentation.slides, 1):
            slide_analysis = self._analyze_slide(slide, idx)
            self.analysis["slides"].append(slide_analysis)

        # Analyse globale
        self._analyze_structure()
        self._analyze_consistency()
        self._generate_summary()

        return self.analysis

    def _analyze_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """Analyse une slide individuelle."""
        analysis = {
            "slide_number": slide_number,
            "title": "",
            "text_boxes": [],
            "word_count": 0,
            "bullet_points": 0,
            "font_sizes": [],
            "issues": [],
            "has_image": False,
            "has_chart": False,
            "has_table": False
        }

        text_content = []

        for shape in slide.shapes:
            # Détection du titre
            if hasattr(shape, "text") and shape.text.strip():
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                    analysis["title"] = shape.text.strip()

                text_content.append(shape.text)
                analysis["text_boxes"].append(shape.text.strip())

                # Comptage des bullets
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level > 0 or paragraph.text.strip().startswith(('•', '-', '*')):
                            analysis["bullet_points"] += 1

                # Extraction des tailles de police
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                analysis["font_sizes"].append(run.font.size.pt)

            # Détection des éléments visuels
            if shape.shape_type == 13:  # Picture
                analysis["has_image"] = True
            elif shape.shape_type == 3:  # Chart
                analysis["has_chart"] = True
            elif shape.shape_type == 19:  # Table
                analysis["has_table"] = True

        # Comptage des mots
        full_text = " ".join(text_content)
        analysis["word_count"] = len(full_text.split())

        # Détection des problèmes
        self._detect_slide_issues(analysis)

        return analysis

    def _detect_slide_issues(self, slide_analysis: Dict[str, Any]) -> None:
        """Détecte les problèmes sur une slide en utilisant les seuils configurables."""
        issues = []

        # Récupérer les seuils depuis la configuration
        max_words = self.config.get('max_words_per_slide', 50)
        max_bullets = self.config.get('max_bullets_per_slide', 7)
        min_words = self.config.get('min_words_per_slide', 5)
        max_font_vars = self.config.get('max_font_variations', 4)

        # Pas de titre
        if not slide_analysis["title"]:
            issues.append({
                "type": "structure",
                "severity": "high",
                "message": "Slide sans titre - chaque slide doit avoir un titre clair"
            })

        # Trop de texte (seuil configurable)
        if slide_analysis["word_count"] > max_words:
            issues.append({
                "type": "clarity",
                "severity": "medium",
                "message": f"Trop de texte ({slide_analysis['word_count']} mots) - limite configurée: {max_words} mots"
            })

        # Trop de bullets (seuil configurable)
        if slide_analysis["bullet_points"] > max_bullets:
            issues.append({
                "type": "clarity",
                "severity": "medium",
                "message": f"Trop de points ({slide_analysis['bullet_points']}) - limite configurée: {max_bullets} points"
            })

        # Incohérence de tailles de police (seuil configurable)
        if len(set(slide_analysis["font_sizes"])) > max_font_vars:
            issues.append({
                "type": "design",
                "severity": "low",
                "message": f"Trop de tailles de police différentes ({len(set(slide_analysis['font_sizes']))}) - limite: {max_font_vars}"
            })

        # Slide trop vide (seuil configurable)
        if slide_analysis["word_count"] < min_words and not any([
            slide_analysis["has_image"],
            slide_analysis["has_chart"],
            slide_analysis["has_table"]
        ]):
            issues.append({
                "type": "structure",
                "severity": "low",
                "message": f"Slide presque vide ({slide_analysis['word_count']} mots) - minimum: {min_words} mots ou ajouter visuel"
            })

        slide_analysis["issues"] = issues

    def _analyze_structure(self) -> None:
        """Analyse la structure globale de la présentation."""
        slides = self.analysis["slides"]

        # Vérifier la présence d'une slide de titre
        if slides and slides[0]["word_count"] < 10:
            pass  # OK, probablement une slide de titre
        else:
            self.analysis["global_issues"].append({
                "type": "structure",
                "severity": "medium",
                "message": "Première slide devrait être une slide de titre concise"
            })

        # Vérifier la présence d'une conclusion
        if slides and "conclu" not in slides[-1]["title"].lower() and "merci" not in slides[-1]["title"].lower():
            self.analysis["global_issues"].append({
                "type": "structure",
                "severity": "low",
                "message": "Dernière slide devrait être une conclusion ou remerciements"
            })

        # Détecter les slides trop denses consécutives (seuils configurables)
        dense_threshold = self.config.get('dense_slide_threshold', 40)
        max_dense_consecutive = self.config.get('max_dense_slides_consecutive', 3)

        dense_slides = [i for i, s in enumerate(slides) if s["word_count"] > dense_threshold]
        if len(dense_slides) > max_dense_consecutive:
            self.analysis["global_issues"].append({
                "type": "clarity",
                "severity": "medium",
                "message": f"{len(dense_slides)} slides sont très denses (>{dense_threshold} mots) - alterner avec des slides visuelles"
            })

    def _analyze_consistency(self) -> None:
        """Analyse la cohérence de la présentation."""
        slides = self.analysis["slides"]

        # Vérifier la cohérence des tailles de police
        all_font_sizes = []
        for slide in slides:
            all_font_sizes.extend(slide["font_sizes"])

        if all_font_sizes:
            font_counter = Counter(all_font_sizes)
            most_common = font_counter.most_common(1)[0][0]

            if len(font_counter) > 6:
                self.analysis["global_issues"].append({
                    "type": "design",
                    "severity": "medium",
                    "message": f"Trop de tailles de police différentes ({len(font_counter)}) - standardiser"
                })

    def _generate_summary(self) -> None:
        """Génère un résumé de l'analyse."""
        slides = self.analysis["slides"]

        total_issues = len(self.analysis["global_issues"])
        for slide in slides:
            total_issues += len(slide["issues"])

        high_severity = sum(1 for slide in slides for issue in slide["issues"] if issue["severity"] == "high")
        high_severity += sum(1 for issue in self.analysis["global_issues"] if issue["severity"] == "high")

        avg_words = sum(s["word_count"] for s in slides) / len(slides) if slides else 0
        avg_bullets = sum(s["bullet_points"] for s in slides) / len(slides) if slides else 0

        self.analysis["summary"] = {
            "total_issues": total_issues,
            "high_severity_issues": high_severity,
            "avg_words_per_slide": round(avg_words, 1),
            "avg_bullets_per_slide": round(avg_bullets, 1),
            "slides_with_images": sum(1 for s in slides if s["has_image"]),
            "slides_with_charts": sum(1 for s in slides if s["has_chart"]),
            "slides_with_tables": sum(1 for s in slides if s["has_table"])
        }

    def print_report(self) -> None:
        """Affiche un rapport d'analyse lisible."""
        print("=" * 80)
        print(f"📊 RAPPORT D'ANALYSE: {self.analysis['filename']}")
        print("=" * 80)
        print()

        summary = self.analysis["summary"]
        print(f"📈 Statistiques globales:")
        print(f"   • Total de slides: {self.analysis['total_slides']}")
        print(f"   • Mots par slide (moyenne): {summary['avg_words_per_slide']}")
        print(f"   • Points par slide (moyenne): {summary['avg_bullets_per_slide']}")
        print(f"   • Slides avec images: {summary['slides_with_images']}")
        print(f"   • Slides avec graphiques: {summary['slides_with_charts']}")
        print(f"   • Slides avec tableaux: {summary['slides_with_tables']}")
        print()

        print(f"⚠️  Problèmes détectés:")
        print(f"   • Total: {summary['total_issues']}")
        print(f"   • Sévérité haute: {summary['high_severity_issues']}")
        print()

        if self.analysis["global_issues"]:
            print("🌐 Problèmes globaux:")
            for issue in self.analysis["global_issues"]:
                emoji = "🔴" if issue["severity"] == "high" else "🟡" if issue["severity"] == "medium" else "🟢"
                print(f"   {emoji} [{issue['type'].upper()}] {issue['message']}")
            print()

        print("📑 Analyse par slide:")
        for slide in self.analysis["slides"]:
            if slide["issues"]:
                title = slide["title"] if slide["title"] else "(Sans titre)"
                print(f"\n   Slide {slide['slide_number']}: {title}")
                for issue in slide["issues"]:
                    emoji = "🔴" if issue["severity"] == "high" else "🟡" if issue["severity"] == "medium" else "🟢"
                    print(f"      {emoji} {issue['message']}")

        print()
        print("=" * 80)


def main():
    parser = argparse.ArgumentParser(
        description="Analyse une présentation PowerPoint selon les standards de conseil"
    )
    parser.add_argument("pptx_file", help="Fichier PowerPoint à analyser")
    parser.add_argument("--output", "-o", help="Fichier JSON de sortie (optionnel)")
    parser.add_argument("--quiet", "-q", action="store_true", help="Pas d'affichage console")
    parser.add_argument("--config", "-c", help="Fichier de configuration JSON personnalisé")
    parser.add_argument("--preset", "-p",
                       choices=['conseil', 'executive', 'technique', 'commercial', 'dsfr-strict'],
                       help="Preset de configuration prédéfini")

    args = parser.parse_args()

    try:
        # Afficher la configuration utilisée si verbose
        if not args.quiet and (args.config or args.preset):
            config_info = f"Configuration: {args.preset or 'default'}"
            if args.config:
                config_info += f" + {args.config}"
            print(f"⚙️  {config_info}\n")

        analyzer = PresentationAnalyzer(
            args.pptx_file,
            config_file=args.config,
            preset=args.preset
        )
        analysis = analyzer.analyze()

        if not args.quiet:
            analyzer.print_report()

        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                json.dump(analysis, f, ensure_ascii=False, indent=2)
            print(f"\n✅ Rapport JSON sauvegardé: {args.output}")

    except FileNotFoundError as e:
        print(f"❌ Erreur: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"❌ Erreur inattendue: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
