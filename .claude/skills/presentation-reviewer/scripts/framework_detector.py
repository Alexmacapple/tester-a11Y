#!/usr/bin/env python3
"""
Détecte le framework de storytelling utilisé dans une présentation PowerPoint
et suggère le framework le plus adapté si nécessaire.

Frameworks supportés:
- AIDA (Attention, Intérêt, Désir, Action)
- PASS (Problème, Agitation, Solution, Situation)
- What/So What/Now What
- SCQA (Situation, Complication, Question, Answer)
- MECE (Mutually Exclusive, Collectively Exhaustive)
- Pyramide de Minto

Usage:
    python framework_detector.py <fichier.pptx> [--suggest]
"""

import sys
import argparse
from pathlib import Path
from typing import Dict, List, Tuple, Any
import re

try:
    from pptx import Presentation
except ImportError:
    print("Erreur: Le module python-pptx n'est pas installé.", file=sys.stderr)
    print("Installez-le avec: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


class FrameworkDetector:
    """Détecte le framework de storytelling d'une présentation."""

    # Mots-clés pour détecter chaque framework
    FRAMEWORK_KEYWORDS = {
        "AIDA": {
            "attention": ["attention", "découvrez", "imaginez", "regardez", "savez-vous"],
            "intérêt": ["intérêt", "pourquoi", "avantages", "bénéfices"],
            "désir": ["désir", "solution", "résultat", "transformation"],
            "action": ["action", "commencez", "agissez", "contactez", "inscrivez"]
        },
        "PASS": {
            "problème": ["problème", "défi", "difficulté", "obstacle", "enjeu"],
            "agitation": ["impact", "conséquence", "risque", "coût", "perte"],
            "solution": ["solution", "résoudre", "réponse", "méthode", "approche"],
            "situation": ["résultat", "bénéfice", "amélioration", "gain"]
        },
        "What/So What/Now What": {
            "what": ["qu'est-ce", "quoi", "contexte", "situation", "état"],
            "so_what": ["pourquoi", "impact", "importance", "signification"],
            "now_what": ["maintenant", "prochaine", "action", "étape", "plan"]
        },
        "SCQA": {
            "situation": ["contexte", "situation", "actuellement", "aujourd'hui"],
            "complication": ["problème", "cependant", "mais", "défi", "obstacle"],
            "question": ["question", "comment", "pourquoi", "quel"],
            "answer": ["réponse", "solution", "proposition", "recommandation"]
        }
    }

    STRUCTURE_PATTERNS = {
        "Pyramide": ["synthèse", "recommandation", "raison", "preuve", "argument"],
        "MECE": ["catégorie", "segment", "type", "exclusif", "exhaustif"]
    }

    def __init__(self, pptx_path: str):
        """Initialise le détecteur avec le chemin du fichier PowerPoint."""
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"Fichier introuvable: {pptx_path}")

        self.presentation = Presentation(str(self.pptx_path))
        self.slides_text = self._extract_all_text()

    def _extract_all_text(self) -> List[str]:
        """Extrait le texte de toutes les slides."""
        slides_text = []

        for slide in self.presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.lower())
            slides_text.append(" ".join(slide_text))

        return slides_text

    def detect_framework(self) -> Dict[str, Any]:
        """Détecte le framework utilisé dans la présentation."""
        scores = {}

        # Calcul des scores pour les frameworks narratifs
        for framework, stages in self.FRAMEWORK_KEYWORDS.items():
            score = self._calculate_framework_score(framework, stages)
            scores[framework] = score

        # Calcul des scores pour les structures
        for structure, keywords in self.STRUCTURE_PATTERNS.items():
            score = self._calculate_structure_score(structure, keywords)
            scores[structure] = score

        # Identifier le framework dominant
        detected = max(scores.items(), key=lambda x: x[1])

        result = {
            "detected_framework": detected[0] if detected[1] > 0.2 else "Aucun framework clair détecté",
            "confidence": round(detected[1], 2),
            "all_scores": {k: round(v, 2) for k, v in sorted(scores.items(), key=lambda x: x[1], reverse=True)},
            "framework_progression": []
        }

        # Analyser la progression du framework détecté
        if detected[1] > 0.2 and detected[0] in self.FRAMEWORK_KEYWORDS:
            result["framework_progression"] = self._analyze_progression(
                detected[0],
                self.FRAMEWORK_KEYWORDS[detected[0]]
            )

        return result

    def _calculate_framework_score(self, framework: str, stages: Dict[str, List[str]]) -> float:
        """Calcule le score de présence d'un framework."""
        total_text = " ".join(self.slides_text)
        stage_scores = []

        for stage, keywords in stages.items():
            keyword_count = sum(total_text.count(kw) for kw in keywords)
            stage_scores.append(min(keyword_count / 3, 1.0))  # Normaliser à 1.0

        # Score = moyenne des stages détectés
        return sum(stage_scores) / len(stage_scores) if stage_scores else 0.0

    def _calculate_structure_score(self, structure: str, keywords: List[str]) -> float:
        """Calcule le score de présence d'une structure."""
        total_text = " ".join(self.slides_text)
        keyword_count = sum(total_text.count(kw) for kw in keywords)

        return min(keyword_count / 5, 1.0)  # Normaliser à 1.0

    def _analyze_progression(self, framework: str, stages: Dict[str, List[str]]) -> List[Dict[str, Any]]:
        """Analyse la progression du framework à travers les slides."""
        progression = []

        for idx, slide_text in enumerate(self.slides_text, 1):
            stage_matches = {}

            for stage, keywords in stages.items():
                matches = sum(slide_text.count(kw) for kw in keywords)
                if matches > 0:
                    stage_matches[stage] = matches

            if stage_matches:
                dominant_stage = max(stage_matches.items(), key=lambda x: x[1])
                progression.append({
                    "slide": idx,
                    "stage": dominant_stage[0],
                    "confidence": min(dominant_stage[1] / 2, 1.0)
                })

        return progression

    def suggest_framework(self, presentation_type: str = "general") -> Dict[str, Any]:
        """Suggère le framework le plus adapté selon le type de présentation."""
        suggestions = {
            "commercial": {
                "primary": "AIDA",
                "reason": "AIDA est idéal pour les présentations commerciales et la vente de produits/services",
                "structure": [
                    "Slide 1-2: Capter l'attention avec un fait marquant ou une question",
                    "Slide 3-4: Susciter l'intérêt en présentant le contexte et les enjeux",
                    "Slide 5-7: Créer le désir en montrant la solution et ses bénéfices",
                    "Slide 8-9: Appel à l'action clair avec prochaines étapes"
                ],
                "alternatives": ["PASS"]
            },
            "problème": {
                "primary": "PASS",
                "reason": "PASS est parfait pour présenter un problème et sa solution",
                "structure": [
                    "Slide 1-2: Exposer clairement le problème",
                    "Slide 3-4: Agiter en montrant les impacts et conséquences",
                    "Slide 5-7: Présenter la solution de manière détaillée",
                    "Slide 8-9: Décrire la situation future améliorée"
                ],
                "alternatives": ["SCQA"]
            },
            "stratégie": {
                "primary": "SCQA",
                "reason": "SCQA est idéal pour les présentations stratégiques et analytiques",
                "structure": [
                    "Slide 1-2: Décrire la situation actuelle",
                    "Slide 3-4: Identifier les complications et défis",
                    "Slide 5: Poser la question clé à résoudre",
                    "Slide 6-9: Apporter la réponse avec recommandations"
                ],
                "alternatives": ["Pyramide"]
            },
            "compte-rendu": {
                "primary": "What/So What/Now What",
                "reason": "Structure claire pour les comptes-rendus et updates",
                "structure": [
                    "Slide 1-3: What - Qu'est-ce qui s'est passé / situation actuelle",
                    "Slide 4-6: So What - Pourquoi c'est important / impact",
                    "Slide 7-9: Now What - Prochaines étapes / plan d'action"
                ],
                "alternatives": ["SCQA"]
            },
            "conseil": {
                "primary": "Pyramide",
                "reason": "Structure privilégiée dans le conseil pour argumenter et convaincre",
                "structure": [
                    "Slide 1: Message clé / recommandation principale",
                    "Slide 2-4: Arguments principaux (niveau 1)",
                    "Slide 5-8: Preuves et données à l'appui (niveau 2)",
                    "Slide 9: Synthèse et prochaines étapes"
                ],
                "alternatives": ["MECE", "SCQA"]
            },
            "general": {
                "primary": "SCQA",
                "reason": "Framework polyvalent adapté à la plupart des contextes",
                "structure": [
                    "Slide 1-2: Situation",
                    "Slide 3-4: Complication",
                    "Slide 5: Question",
                    "Slide 6-9: Answer"
                ],
                "alternatives": ["What/So What/Now What"]
            }
        }

        return suggestions.get(presentation_type, suggestions["general"])

    def print_report(self, detection: Dict[str, Any]) -> None:
        """Affiche un rapport de détection."""
        print("=" * 80)
        print(f"🎯 DÉTECTION DE FRAMEWORK: {self.pptx_path.name}")
        print("=" * 80)
        print()

        print(f"📊 Framework détecté: {detection['detected_framework']}")
        print(f"   Confiance: {detection['confidence'] * 100:.0f}%")
        print()

        if detection["all_scores"]:
            print("📈 Scores de tous les frameworks:")
            for framework, score in detection["all_scores"].items():
                bar_length = int(score * 20)
                bar = "█" * bar_length + "░" * (20 - bar_length)
                print(f"   {framework:20s} {bar} {score * 100:5.1f}%")
            print()

        if detection["framework_progression"]:
            print("📑 Progression du framework à travers les slides:")
            for prog in detection["framework_progression"]:
                print(f"   Slide {prog['slide']:2d}: {prog['stage']:15s} (confiance: {prog['confidence'] * 100:.0f}%)")
            print()

        print("=" * 80)


def main():
    parser = argparse.ArgumentParser(
        description="Détecte le framework de storytelling d'une présentation PowerPoint"
    )
    parser.add_argument("pptx_file", help="Fichier PowerPoint à analyser")
    parser.add_argument(
        "--suggest",
        choices=["commercial", "problème", "stratégie", "compte-rendu", "conseil", "general"],
        help="Suggérer un framework adapté au type de présentation"
    )

    args = parser.parse_args()

    try:
        detector = FrameworkDetector(args.pptx_file)
        detection = detector.detect_framework()

        detector.print_report(detection)

        if args.suggest:
            print("\n💡 SUGGESTION DE FRAMEWORK")
            print("=" * 80)
            suggestion = detector.suggest_framework(args.suggest)

            print(f"\n📌 Framework recommandé: {suggestion['primary']}")
            print(f"   Raison: {suggestion['reason']}")
            print()
            print("📋 Structure recommandée:")
            for step in suggestion['structure']:
                print(f"   • {step}")
            print()
            print(f"🔄 Alternatives: {', '.join(suggestion['alternatives'])}")
            print()

    except FileNotFoundError as e:
        print(f"❌ Erreur: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"❌ Erreur inattendue: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
