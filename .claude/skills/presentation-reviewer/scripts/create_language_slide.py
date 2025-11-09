#!/usr/bin/env python3
"""
Créer une slide PowerPoint conforme DSFR sur l'identification de la langue
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Couleurs DSFR
BLEU_FRANCE = RGBColor(0, 0, 145)      # #000091
BLEU_INFO = RGBColor(0, 99, 203)       # #0063CB
ROUGE_MARIANNE = RGBColor(225, 0, 15)  # #E1000F
GRIS_1000 = RGBColor(22, 22, 22)       # #161616
BLANC = RGBColor(255, 255, 255)        # #FFFFFF
GRIS_FOND = RGBColor(246, 246, 246)    # #F6F6F6

def create_presentation():
    """Créer la présentation avec la slide sur la langue"""

    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9
    prs.slide_height = Inches(5.625)

    # Ajouter une slide vierge
    blank_slide_layout = prs.slide_layouts[6]  # Layout vide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLANC

    # 1. TITRE
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Identifier la langue principale et les changements de langue"
    p.font.name = "Arial"  # Marianne non disponible par défaut
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLEU_FRANCE

    # 2. LIGNE DE SÉPARATION 1
    line1 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1.0), Inches(9), Inches(0)
    )
    line1.line.color.rgb = GRIS_1000
    line1.line.width = Pt(1)

    # 3. IMPACT UTILISATEURS
    impact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.7)
    )
    impact_frame = impact_box.text_frame
    impact_frame.word_wrap = True
    p = impact_frame.paragraphs[0]
    p.text = "Les personnes aveugles ou malvoyantes qui utilisent un lecteur d'écran peuvent rencontrer d'importantes difficultés à comprendre les textes lus dans une langue qui n'est pas la bonne. Sans indication de langue, le lecteur d'écran utilisera sa langue par défaut, rendant le contenu incompréhensible."
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.color.rgb = GRIS_1000
    p.line_spacing = 1.2

    # 4. LIGNE DE SÉPARATION 2
    line2 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(2.0), Inches(9), Inches(0)
    )
    line2.line.color.rgb = GRIS_1000
    line2.line.width = Pt(1)

    # 5. MESSAGE CLÉ (Bandeau bleu)
    message_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(2.2), Inches(9), Inches(0.6)
    )
    message_box.fill.solid()
    message_box.fill.fore_color.rgb = BLEU_INFO
    message_box.line.color.rgb = BLEU_INFO

    message_frame = message_box.text_frame
    message_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    message_frame.word_wrap = True
    p = message_frame.paragraphs[0]
    p.text = "Sans langue, ou avec une mauvaise indication de langue, pas de lecture accessible"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLANC

    # 6. LIGNE DE SÉPARATION 3
    line3 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(2.9), Inches(9), Inches(0)
    )
    line3.line.color.rgb = GRIS_1000
    line3.line.width = Pt(1)

    # 7. CRITÈRES RGAA
    criteres_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.1), Inches(9), Inches(0.3)
    )
    p = criteres_title.text_frame.paragraphs[0]
    p.text = "Critères RGAA : 8.3, 8.4, 8.7, 8.8, 8.10"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRIS_1000

    # 8. SECTION CODE - 2 colonnes
    # Colonne gauche : Langue principale
    col1_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(4.3), Inches(0.8)
    )
    col1_frame = col1_box.text_frame
    col1_frame.word_wrap = True

    p1 = col1_frame.paragraphs[0]
    p1.text = "Langue principale :"
    p1.font.name = "Arial"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = GRIS_1000

    p2 = col1_frame.add_paragraph()
    p2.text = '<html lang="fr">'
    p2.font.name = "Courier New"
    p2.font.size = Pt(13)
    p2.font.color.rgb = ROUGE_MARIANNE
    p2.space_before = Pt(6)

    # Colonne droite : Changement de langue
    col2_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(3.5), Inches(4.3), Inches(0.8)
    )
    col2_frame = col2_box.text_frame
    col2_frame.word_wrap = True

    p3 = col2_frame.paragraphs[0]
    p3.text = "Changement de langue :"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = GRIS_1000

    p4 = col2_frame.add_paragraph()
    p4.text = '<span lang="en">Hello</span>'
    p4.font.name = "Courier New"
    p4.font.size = Pt(13)
    p4.font.color.rgb = ROUGE_MARIANNE
    p4.space_before = Pt(6)

    # 9. LIGNE DE SÉPARATION 4
    line4 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(4.5), Inches(9), Inches(0)
    )
    line4.line.color.rgb = GRIS_1000
    line4.line.width = Pt(1)

    # 10. OUTILS DE VÉRIFICATION
    outils_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.7), Inches(9), Inches(0.3)
    )
    p = outils_title.text_frame.paragraphs[0]
    p.text = "✅ 3 outils de vérification"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRIS_1000

    # 11. LISTE OUTILS
    outils_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.0), Inches(9), Inches(0.5)
    )
    outils_frame = outils_box.text_frame
    outils_frame.word_wrap = True

    outils = [
        "• Validateur HTML du W3C",
        "• Console du navigateur (F12)",
        "• Lecteurs d'écran (NVDA, JAWS, VoiceOver)"
    ]

    for i, outil in enumerate(outils):
        if i == 0:
            p = outils_frame.paragraphs[0]
        else:
            p = outils_frame.add_paragraph()
        p.text = outil
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = GRIS_1000
        p.space_after = Pt(3)

    return prs

def main():
    """Fonction principale"""
    prs = create_presentation()
    output_file = "slide-langue-dsfr.pptx"
    prs.save(output_file)
    print(f"✅ Slide créée : {output_file}")
    print(f"📊 Caractéristiques :")
    print(f"   - Format : 16:9")
    print(f"   - Couleurs : DSFR (Bleu France, Rouge Marianne)")
    print(f"   - Typographie : Arial (remplacer par Marianne)")
    print(f"   - Critères RGAA : 8.3, 8.4, 8.7, 8.8, 8.10")
    print(f"   - Impact utilisateurs : inclus")
    print(f"   - Titre : Identifier la langue principale et les changements de langue")

if __name__ == "__main__":
    main()
