#!/usr/bin/env python3
"""
Créer une slide PowerPoint conforme DSFR sur les contrastes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Couleurs DSFR
BLEU_FRANCE = RGBColor(0, 0, 145)      # #000091
BLEU_INFO = RGBColor(0, 99, 203)       # #0063CB
ROUGE_MARIANNE = RGBColor(225, 0, 15)  # #E1000F
GRIS_1000 = RGBColor(22, 22, 22)       # #161616
BLANC = RGBColor(255, 255, 255)        # #FFFFFF
GRIS_FOND = RGBColor(246, 246, 246)    # #F6F6F6

def create_presentation():
    """Créer la présentation avec la slide sur les contrastes"""

    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9
    prs.slide_height = Inches(5.625)

    # Ajouter une slide vierge
    blank_slide_layout = prs.slide_layouts[6]  # Layout vide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLANC

    # 1. TITRE
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Appliquer des contrastes suffisamment élevés"
    p.font.name = "Arial"  # Marianne non disponible par défaut
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLEU_FRANCE

    # 2. LIGNE DE SÉPARATION 1
    line1 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1.0), Inches(9), Inches(0)
    )
    line1.line.color.rgb = GRIS_1000
    line1.line.width = Pt(1)

    # 3. IMPACT UTILISATEURS
    impact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.8)
    )
    impact_frame = impact_box.text_frame
    impact_frame.word_wrap = True
    p = impact_frame.paragraphs[0]
    p.text = "L'utilisation de contrastes trop faibles diminue considérablement la lisibilité des contenus. Dans ce cas, les personnes malvoyantes ne perçoivent pas certains textes, composants d'interface ou éléments graphiques. Elles sont alors incapables de lire le contenu ou d'interagir avec les éléments concernés."
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = GRIS_1000
    p.line_spacing = 1.2

    # 4. LIGNE DE SÉPARATION 2
    line2 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(2.1), Inches(9), Inches(0)
    )
    line2.line.color.rgb = GRIS_1000
    line2.line.width = Pt(1)

    # 5. MESSAGE CLÉ (Bandeau bleu)
    message_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(2.3), Inches(9), Inches(0.6)
    )
    message_box.fill.solid()
    message_box.fill.fore_color.rgb = BLEU_INFO
    message_box.line.color.rgb = BLEU_INFO

    message_frame = message_box.text_frame
    message_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    message_frame.word_wrap = True
    p = message_frame.paragraphs[0]
    p.text = "💡 Le contraste : pas de l'esthétique, de l'accessibilité"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLANC

    # 6. LIGNE DE SÉPARATION 3
    line3 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(3.0), Inches(9), Inches(0)
    )
    line3.line.color.rgb = GRIS_1000
    line3.line.width = Pt(1)

    # 7. SECTION RATIOS
    ratio_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), Inches(9), Inches(0.3)
    )
    p = ratio_title.text_frame.paragraphs[0]
    p.text = "Ratios WCAG requis"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRIS_1000

    # Colonne 1 : Texte normal
    col1_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.6), Inches(2.8), Inches(0.9)
    )
    col1_frame = col1_box.text_frame
    col1_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Titre colonne 1
    p1 = col1_frame.paragraphs[0]
    p1.text = "Texte normal"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Arial"
    p1.font.size = Pt(14)
    p1.font.color.rgb = GRIS_1000

    # Description colonne 1
    p1b = col1_frame.add_paragraph()
    p1b.text = "(< 18 pt ou < 14 pt gras)"
    p1b.alignment = PP_ALIGN.CENTER
    p1b.font.name = "Arial"
    p1b.font.size = Pt(11)
    p1b.font.color.rgb = GRIS_1000

    # Ratio colonne 1
    p2 = col1_frame.add_paragraph()
    p2.text = "4,5:1"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Arial"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = ROUGE_MARIANNE

    # Colonne 2 : Grand texte
    col2_box = slide.shapes.add_textbox(
        Inches(3.6), Inches(3.6), Inches(2.8), Inches(0.9)
    )
    col2_frame = col2_box.text_frame
    col2_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Titre colonne 2
    p3 = col2_frame.paragraphs[0]
    p3.text = "Grand texte"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRIS_1000

    # Description colonne 2
    p3b = col2_frame.add_paragraph()
    p3b.text = "(≥ 18 pt ou ≥ 14 pt gras)"
    p3b.alignment = PP_ALIGN.CENTER
    p3b.font.name = "Arial"
    p3b.font.size = Pt(11)
    p3b.font.color.rgb = GRIS_1000

    # Ratio colonne 2
    p4 = col2_frame.add_paragraph()
    p4.text = "3:1"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = "Arial"
    p4.font.size = Pt(32)
    p4.font.bold = True
    p4.font.color.rgb = ROUGE_MARIANNE

    # Colonne 3 : Icônes informatives
    col3_box = slide.shapes.add_textbox(
        Inches(6.7), Inches(3.6), Inches(2.8), Inches(0.9)
    )
    col3_frame = col3_box.text_frame
    col3_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Titre colonne 3
    p5 = col3_frame.paragraphs[0]
    p5.text = "Icônes informatives"
    p5.alignment = PP_ALIGN.CENTER
    p5.font.name = "Arial"
    p5.font.size = Pt(14)
    p5.font.color.rgb = GRIS_1000

    # Description colonne 3
    p5b = col3_frame.add_paragraph()
    p5b.text = "(composants d'interface)"
    p5b.alignment = PP_ALIGN.CENTER
    p5b.font.name = "Arial"
    p5b.font.size = Pt(11)
    p5b.font.color.rgb = GRIS_1000

    # Ratio colonne 3
    p6 = col3_frame.add_paragraph()
    p6.text = "3:1"
    p6.alignment = PP_ALIGN.CENTER
    p6.font.name = "Arial"
    p6.font.size = Pt(32)
    p6.font.bold = True
    p6.font.color.rgb = ROUGE_MARIANNE

    # 8. LIGNE DE SÉPARATION 4
    line4 = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(4.6), Inches(9), Inches(0)
    )
    line4.line.color.rgb = GRIS_1000
    line4.line.width = Pt(1)

    # 9. SOLUTION - Titre
    solution_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.8), Inches(9), Inches(0.3)
    )
    p = solution_title.text_frame.paragraphs[0]
    p.text = "✅ 3 outils de vérification"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRIS_1000

    # 10. LISTE OUTILS
    outils_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.15), Inches(9), Inches(0.4)
    )
    outils_frame = outils_box.text_frame
    outils_frame.word_wrap = True

    outils = [
        "• Color Contrast Analyzer (CCA) - TPGi",
        "• Contrast Finder - Tanaguru",
        "• Analyseur de contraste - Deque University"
    ]

    for i, outil in enumerate(outils):
        if i == 0:
            p = outils_frame.paragraphs[0]
        else:
            p = outils_frame.add_paragraph()
        p.text = outil
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = GRIS_1000
        p.space_after = Pt(3)

    return prs

def main():
    """Fonction principale"""
    prs = create_presentation()
    output_file = "slide-contrastes-dsfr.pptx"
    prs.save(output_file)
    print(f"✅ Slide créée : {output_file}")
    print(f"📊 Caractéristiques :")
    print(f"   - Format : 16:9")
    print(f"   - Couleurs : DSFR (Bleu France, Rouge Marianne)")
    print(f"   - Typographie : Arial (remplacer par Marianne)")
    print(f"   - Structure : 3 colonnes de ratios")
    print(f"   - Impact utilisateurs : inclus")
    print(f"   - Titre : Appliquer des contrastes suffisamment élevés")

if __name__ == "__main__":
    main()
